import io
from typing import Dict, Any
import logging

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

class PPTXGenerator:
    @staticmethod
    def generate(title: str, analytics: Dict[str, Any], forecast: Dict[str, Any], recommendations: Dict[str, Any]) -> bytes:
        if not HAS_PPTX:
            # Fallback simple text formatting for pptx binary stream mock
            logger.warning("python-pptx is not installed. Generating simulated PowerPoint pptx binary stream.")
            buffer = io.BytesIO()
            buffer.write(f"=== PPTX Presentation: {title} ===\n".encode("utf-8"))
            return buffer.getvalue()

        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Business Intelligence Audit Presentation"
        
        # KPIs Slide
        slide_layout2 = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(slide_layout2)
        slide2.shapes.title.text = "KPI Performance Summary"
        tf2 = slide2.placeholders[1].text_frame
        for k, v in list(analytics.get("kpis", {}).items())[:5]:
            p = tf2.add_paragraph()
            p.text = f"{k}: {v}"

        # Forecast Slide
        slide3 = prs.slides.add_slide(slide_layout2)
        slide3.shapes.title.text = "Forecasting Metrics Trends"
        tf3 = slide3.placeholders[1].text_frame
        tf3.add_paragraph().text = f"Trend Direction: {forecast.get('trend', 'Flat')}"
        tf3.add_paragraph().text = f"Model Selected: {forecast.get('model_used', 'ARIMA')}"

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
